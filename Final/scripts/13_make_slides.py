"""
13 - Final Presentation slide deck (.pptx) + speaker notes.

14 slides in the same 4-theme structure the Midterm deck used, but
focused on the Final's three predictive paths (link prediction, return
prediction, sentiment fine-tune).

Speaker notes (Thai) are attached to every slide via
    slide.notes_slide.notes_text_frame
so PowerPoint's Presenter View shows them automatically.

Output: report/DADS7201_Final_Slides.pptx
"""
from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

MIDTERM_DIR = Path('../Midterm')
MID_OUT     = MIDTERM_DIR / 'output'
FIN_OUT     = Path('output')
REPORT_DIR  = Path('report')
REPORT_DIR.mkdir(exist_ok=True)

# ----- data ---------------------------------------------------------------
lp = pd.read_csv(FIN_OUT / 'lp_baselines.csv')
lp_topk = pd.read_csv(FIN_OUT / 'lp_topk_predicted_edges.csv')
lp_split = json.load(open(Path('data') / 'temporal_split.json', encoding='utf-8'))
rp = json.load(open(FIN_OUT / 'rp_metrics.json', encoding='utf-8'))
val = json.load(open(FIN_OUT / 'sentiment_finetuned_validation.json', encoding='utf-8'))

# Midterm context
evo = pd.read_csv(MID_OUT / 'evolution_comparison.csv').set_index('Phase')
sig = json.load(open(MID_OUT / 'significance.json', encoding='utf-8'))
overlap = pd.read_csv(MID_OUT / 'overlap_sensitivity.csv').sort_values('Jaccard', ascending=False)
primary = overlap.iloc[0]

lp_test = lp[lp['split'] == 'test'].sort_values('auc', ascending=False)
best_lp = lp_test.iloc[0]
b0, b1, b2, b3 = rp['B0'], rp['B1'], rp['B2'], rp['B3']

# ----- style helpers ------------------------------------------------------
PAL = {
    'navy':   RGBColor(0x1A, 0x1A, 0x2E),
    'red':    RGBColor(0xC1, 0x27, 0x2D),
    'purple': RGBColor(0x9D, 0x4E, 0xDD),
    'grey':   RGBColor(0x66, 0x66, 0x66),
    'green':  RGBColor(0x2A, 0x9D, 0x8F),
    'blue':   RGBColor(0x1D, 0x35, 0x57),
    'amber':  RGBColor(0xE2, 0x7D, 0x60),
}
THAI = 'Tahoma'


def _run(run, size, color='navy', bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = PAL[color]
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = THAI


def title_of(s, text, size=28, color='navy'):
    t = s.shapes.title
    t.text = text
    for r in t.text_frame.paragraphs[0].runs:
        _run(r, size, color=color, bold=True)


def tb(s, left, top, w, h, lines, size=14, color='navy', bullet=False,
       italic=False):
    box = s.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ('• ' if bullet else '') + line
        for r in p.runs:
            _run(r, size, color=color, italic=italic)


def img(s, name, left, top, *, width=None, height=None):
    for src in (FIN_OUT, MID_OUT):
        p = src / name
        if p.exists():
            if width:
                s.shapes.add_picture(str(p), left, top, width=width)
            else:
                s.shapes.add_picture(str(p), left, top, height=height)
            return


def chip(s, label, color='blue'):
    box = s.shapes.add_textbox(Inches(0.35), Inches(0.25), Inches(4), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = label
    for r in p.runs:
        _run(r, 12, color=color, bold=True)


def notes(s, text):
    """Attach speaker notes to a slide."""
    tf = s.notes_slide.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = 'Tahoma'
            r.font.size = Pt(12)


# ========================================================================
# Deck
# ========================================================================
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ----- Slide 1: Title ---------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
title_of(s, 'Predictive Extension of the SET50 Hype Network', size=32)
tb(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.0),
   ['Temporal Link Prediction, Return Forecasting, and Sentiment Fine-tuning'],
   size=20, color='grey')
tb(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.5),
   [
       'DADS7201 — Final Project · NIDA',
       'Puwadon Sri · กรกฎาคม 2026',
       ('ต่อยอด Midterm SNA pipeline ด้วย predictive analytics '
        '3 ทาง (Weeks 3–7 techniques)'),
   ],
   size=16)
notes(s, (
    'สวัสดีครับอาจารย์และเพื่อนๆ วันนี้ผมจะนำเสนอ Final Project ที่ต่อยอด '
    'จาก Midterm ครับ\n\n'
    'ตอน Midterm เราสร้าง SNA pipeline บน SET50 + Pantip Sinthorn — เป็นงาน '
    'descriptive analytics คือ อธิบายว่ามีอะไรเกิดขึ้น\n\n'
    'Final ครั้งนี้ต่อด้วย predictive analytics 3 ทาง (Link Prediction, '
    'Return Prediction, Sentiment Fine-tune) โดยใช้เทคนิคจาก Week 3 ถึง Week 7 '
    'ที่เรียนมาทั้งหมด\n\n'
    'ใช้เวลาประมาณ 15 นาที มีคำถามทักได้เลยครับ'
))


# ----- Slide 2: Recap Midterm -------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'ที่มา · Recap Midterm')
title_of(s, 'Midterm ทำอะไรไว้ (Descriptive)')
tb(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
   [
       ('Pipeline: Pantip Sinthorn scrape (Playwright) → SET50 mention '
        'regex → co-mention graph → Louvain communities → WangchanBERTa '
        'sentiment → Pearson return-correlation overlay'),
       (f'ผลหลัก 4 hypotheses (H1–H4): H1 (degree→attention) ✅, H2 '
        f'(community→co-movement) ✅ p=0.001, H3 (sentiment→hub) ❌ อ่อน, '
        f'H4 (event→density) ✅ p={sig["permutation"]["p_value_after_before"]:.3f}'),
       (f'Attention vs Fundamentals overlay: Jaccard {primary["Jaccard"]:.2f} '
        f'(95% CI [{sig["bootstrap_jaccard"]["ci_95_low"]:.2f}, '
        f'{sig["bootstrap_jaccard"]["ci_95_high"]:.2f}]) — retail attention '
        'กับ fundamentals แทบ disjoint กัน'),
       ('Sentiment classifier ที่ใช้ off-the-shelf: accuracy 0.34, '
        'Cohen κ 0.14 — chance level (จุดอ่อนใหญ่)'),
       ('อาจารย์ Midterm แนะให้ทำ predictive layer + validate classifier '
        '→ ผลักไปทำใน Final'),
   ],
   size=14, bullet=True)
notes(s, (
    'Slide นี้ recap สิ่งที่ทำใน Midterm สั้นๆ เพราะ Final ต่อยอดจากตรงนี้\n\n'
    'Pipeline: เราเก็บโพสต์จริงจาก Pantip ห้องสินธร ด้วย Playwright, '
    'สกัด SET50 mentions, สร้าง co-mention graph, run Louvain, '
    'ประเมิน sentiment ด้วย WangchanBERTa, และ overlay กับ correlation network\n\n'
    'ผลหลัก: H1, H2, H4 supported (มี p-value รองรับ) แต่ H3 อ่อน\n\n'
    'จุดเด่นคือ Jaccard 0.14 — เราเจอว่า retail attention กับ fundamentals '
    'ไม่ค่อยไปทางเดียวกัน\n\n'
    'จุดอ่อนใหญ่: classifier chance-level อาจารย์เลยแนะให้ต่อยอดตรงนี้'
))


# ----- Slide 3: Predictive Framing --------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'ที่มา · Predictive Framing')
title_of(s, 'จาก Descriptive → Predictive')
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.7),
   ['ทำ Midterm finding ให้กลายเป็น testable prediction'],
   size=16, color='grey', italic=True)
tb(s, Inches(0.7), Inches(2.3), Inches(12), Inches(5),
   [
       'H4 (event→density) → Link Prediction — ทำนายได้ไหมว่า pair ไหน '
       'จะ co-mention กันในอนาคต?',
       'Overlay (attention ≠ fundamentals) → Predictive test — model จะไป '
       'concentrate ที่ hype-only หรือ overlooked pairs?',
       'Hype Hub Score + graph features → Return Prediction — เพิ่ม signal '
       'เกินกว่า own-return autocorrelation ไหม?',
       'H3 อ่อน → Classifier fix — เป็น data property หรือ classifier '
       'artefact?',
   ],
   size=16, bullet=True)
tb(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.7),
   ['ผลลัพธ์ทุกด้าน = falsifiable ไม่ว่าตอบ yes หรือ no ก็เขียนได้'],
   size=13, italic=True, color='grey')
notes(s, (
    'Slide สำคัญ — อธิบาย logic ว่าทำไมเราเลือก 3 tasks นี้\n\n'
    'ทุก finding descriptive ของ Midterm มี testable prediction แฝงอยู่:\n'
    '- H4 density เพิ่ม → ทำนายได้ไหมว่า edge ใหม่จะเกิดที่ไหน?\n'
    '- Overlay Jaccard ต่ำ → model จะ predict hype-only หรือ overlooked?\n'
    '- Hype Hub → ใช้ทำนาย return ได้ไหม?\n'
    '- Classifier chance-level → fix ได้ไหม?\n\n'
    'จุดสำคัญ: ผลลัพธ์ทุกด้าน falsifiable — ไม่ว่าตอบ yes หรือ no ก็เป็น '
    'academic contribution ได้'
))


# ----- Slide 4: Overview of 3 Paths -------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'การดำเนินการ · 3 Paths', color='green')
title_of(s, '3 Predictive Paths (Weeks 3–7 techniques)')
tb(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
   [
       ('🥇 Path A — Temporal Link Prediction (Weeks 3, 5, 7)'),
       ('   ทำนายว่า SET50 pair ไหนจะ co-mention กันครั้งแรกใน 20 วันข้างหน้า'),
       ('   เทียบ 5 methods: Jaccard, Adamic-Adar, PA, Node2Vec, GraphSAGE'),
       (''),
       ('🥈 Path B — Return Prediction (finance grounding)'),
       ('   ทำนาย sign(next-day return) ใช้ graph + sentiment features'),
       ('   B0 (lag only) → B1 → B2 → B3 (XGBoost), 5-fold expanding CV'),
       (''),
       ('🥉 Path C — Sentiment Classifier Fine-tune (Midterm limitation)'),
       ('   Fine-tune WangchanBERTa บน 50 gold posts, 5-fold CV, 3 epochs each'),
       ('   วัด accuracy / F1 / Cohen κ before vs after'),
   ],
   size=14)
notes(s, (
    '3 paths ที่เลือกทำ:\n\n'
    'Path A — Link Prediction: ต่อยอด Week 7 workshop โดยตรง เอา heuristics '
    'จาก Week 3 (Jaccard, Adamic-Adar, PA) มาเทียบกับ Node2Vec (Week 5) และ '
    'GraphSAGE (Week 7) — ครอบคลุมทั้ง semester\n\n'
    'Path B — Return Prediction: เอา graph features จาก Midterm มา predict '
    'next-day return sign — ตอบโจทย์ที่อาจารย์แนะไว้\n\n'
    'Path C — Sentiment Fine-tune: แก้จุดอ่อน chance-level classifier ที่เจอ '
    'ใน Midterm ด้วย 5-fold CV บน 50 gold posts\n\n'
    'ทั้ง 3 paths ใช้เทคนิคที่เรียนใน Week 3 ถึง Week 7 ครบทุก week'
))


# ----- Slide 5: Path A methodology ---------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'การดำเนินการ · Path A Methodology', color='green')
title_of(s, 'Path A — Temporal Link Prediction Setup')
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(3.5),
   [
       (f'Temporal split (event-anchored T = {lp_split["event_date"]}):'),
       (f'   Train = co-mentions in [{lp_split["train_end"]} back]'
        f'  ({lp_split["n_train_edges"]} edges)'),
       (f'   Val = first-time edges in [{lp_split["train_end"]} → '
        f'{lp_split["val_end"]}]  ({lp_split["n_val_positive_first_seen"]} positives)'),
       (f'   Test = first-time edges in [{lp_split["val_end"]} → '
        f'{lp_split["test_end"]}]  ({lp_split["n_test_positive_first_seen"]} positives)'),
       ('Negatives sampled 2:1 (train), 1:1 (eval) จากคู่ที่ไม่มีในกราฟ'),
   ],
   size=14, bullet=True)
tb(s, Inches(0.7), Inches(4.7), Inches(12), Inches(2.5),
   [
       '5 methods เปรียบเทียบ:',
       ('   Week 3 heuristics — Jaccard, Adamic-Adar, Preferential Attachment'),
       ('   Week 5 — Node2Vec (dim=64, walks=15, walk_len=40) + '
        'LogReg บน Hadamard z_u ⊙ z_v'),
       ('   Week 7 — GraphSAGE (2 SAGEConv layers, hidden=64, '
        'BCE + Adam 5e-3, 250 epochs)'),
       ('Metrics: AUC + Precision@{10,25,50} + MRR + bootstrap 95% CI (500 iters)'),
   ],
   size=14, bullet=True)
notes(s, (
    'Path A methodology — สำคัญมาก\n\n'
    'Temporal split: ต่างจาก random split ตรงที่เราใช้ time เป็นตัวแบ่ง — '
    'train ก่อน T-20, val [T-20, T), test [T, T+20]\n\n'
    'สำหรับ positive edges ใน val/test เราเลือกเฉพาะ "first-time" pairs — '
    'edges ที่ไม่เคยเห็นใน train เท่านั้น (ป้องกัน leakage)\n\n'
    '5 methods เลือกเพื่อครอบคลุมทั้ง 3 tiers ของ complexity:\n'
    '- Heuristics (no training) — ควรเป็น baseline\n'
    '- Shallow embeddings (Week 5 Node2Vec) — เรียนโครงสร้างกราฟ\n'
    '- Deep GNN (Week 7 GraphSAGE) — มี inductive bias\n\n'
    'Bootstrap CI สำคัญเพราะเราต้องมี statistical rigor — Midterm feedback '
    'ระบุจุดนี้'
))


# ----- Slide 6: Path A Results ------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'สรุปผล · Path A Results', color='red')
title_of(s, 'Path A ผลลัพธ์ — Heuristics ชนะ GNN!')
rows_txt = ['Method            AUC     95% CI          P@10   P@25   MRR']
for _, r in lp_test.iterrows():
    rows_txt.append(
        f'{r["method"]:16s}  {r["auc"]:.3f}  '
        f'[{r["auc_ci_low"]:.2f}, {r["auc_ci_high"]:.2f}]  '
        f'{r["precision_at_10"]:.2f}    {r["precision_at_25"]:.2f}   '
        f'{r["mrr"]:.3f}'
    )
tb(s, Inches(0.7), Inches(1.4), Inches(7.5), Inches(2.5), rows_txt, size=13)
img(s, 'lp_roc.png', Inches(8.4), Inches(1.4), width=Inches(4.7))
tb(s, Inches(0.7), Inches(4.6), Inches(12), Inches(2.7),
   [
       (f'Best: {best_lp["method"]} — AUC {best_lp["auc"]:.3f} '
        f'(P@10 = {best_lp["precision_at_10"]:.2f} — 10/10!)'),
       ('Node2Vec (Week 5) เกาะติดกับ heuristics ที่ ~0.71'),
       ('GraphSAGE (Week 7) AUC เพียง 0.55 — ตกลงมาเพราะ:'),
       ('   (1) กราฟเล็ก 50 nodes / 154 edges — training data ไม่พอ'),
       ('   (2) identity features → ไม่มี prior ให้ SAGEConv เรียนรู้'),
       ('Finding เข้ากับ literature — Lü & Zhou (2011) survey ระบุว่า simple '
        'heuristics มัก competitive กับ deep methods ที่ small sparse graphs'),
   ],
   size=13, bullet=True)
notes(s, (
    'Path A ผลลัพธ์ — findings คลาสสิก!\n\n'
    'Preferential Attachment ชนะที่ AUC 0.733 P@10 = 1.00 (10/10 top predictions '
    'ถูกทั้งหมด!) เพราะ retail attention เป็น self-reinforcing — คนพูดถึงหุ้น '
    'ดัง (high degree) มากกว่าหุ้นเงียบ\n\n'
    'GraphSAGE แย่กว่าเยอะ — AUC 0.55 เพราะ:\n'
    '- 50 nodes / 154 training edges เล็กเกินไปสำหรับ GNN\n'
    '- identity features (one-hot node id) → SAGEConv ไม่มีอะไรให้เรียน\n\n'
    'ผลนี้จริงๆ เป็น "positive finding" ในเชิง academic — ตรงกับ literature '
    'ว่า heuristics competitive กับ GNN ที่กราฟเล็ก\n\n'
    'เตรียม defense: อาจารย์อาจถามว่าทำไม GNN แพ้ ให้ตอบว่าเพราะกราฟเล็ก + '
    'ไม่มี node features — ถ้าใส่ features เช่น sector one-hot หรือ centrality '
    'ก็อาจดีขึ้น'
))


# ----- Slide 7: Top Predicted Pairs -------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'สรุปผล · Path A Story', color='red')
title_of(s, 'ผลจริง — Retail แพร่ Hype ไปที่หุ้นดังอีก')
top10 = lp_topk.head(10)
rows_txt = ['Rank  Pair                Score   Midterm bucket        Emerged?']
for i, (_, r) in enumerate(top10.iterrows(), 1):
    rows_txt.append(
        f'  {i:2d}   {r["StockA"]}–{r["StockB"]:<10s}  {r["PredScore"]:>4.0f}   '
        f'{r["MidtermBucket"]:20s}  {r["EmergedInEval"]}'
    )
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(4.5),
   rows_txt, size=13)
tb(s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.3),
   [
       ('Top-10 predicted pairs ส่วนมากเป็น "new" หรือ "hype-only" bucket'),
       ('Overlooked co-movers (จาก Midterm — เช่น CPF-HMPRO, EGCO-RATCH) '
        'ไม่ค่อยติด top ranking → retail ยังไม่ "catch up" กับ fundamentals'),
   ],
   size=13, bullet=True, color='red')
notes(s, (
    'Slide นี้เชื่อม Midterm finding เข้ากับ Path A ผลลัพธ์\n\n'
    'Model predict top-10 unseen pairs → ดูว่าตรงกับ Midterm bucket ไหน:\n'
    '- ถ้า overlap เยอะกับ "overlooked co-movers" → retail กำลังจะ catch up\n'
    '- ถ้า overlap เยอะกับ "hype-only" หรือ "new" → retail ยัง self-reinforce\n\n'
    'ผลจริง: ส่วนมากเป็น hype-only + new — ยืนยันว่า retail attention '
    'ไม่ได้ diffuse ไปยัง overlooked pairs\n\n'
    'ผลนี้สำคัญมาก — descriptive finding "Jaccard 0.14" ตอน Midterm '
    'ที่บอกว่า attention ≠ fundamentals ตอนนี้กลายเป็น testable prediction '
    'ที่ confirmed แล้ว'
))


# ----- Slide 8: Path B methodology --------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'การดำเนินการ · Path B Methodology', color='green')
title_of(s, 'Path B — Return Prediction Setup')
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(2.5),
   [
       ('Task: predict sign(next-day return) ของ SET50 stocks'),
       ('Data: 120 trading days × 50 stocks จาก yfinance (~5,700 rows)'),
       ('CV: expanding-window time-series, 5 folds'),
       ('Portfolio: long top-decile / short bottom-decile, daily rebalance'),
   ],
   size=15, bullet=True)
tb(s, Inches(0.7), Inches(4.2), Inches(12), Inches(3),
   [
       'Feature groups (4 models):',
       ('   B0 (baseline) — own-return lags: r_lag1, r_lag5 mean/std, '
        'cross-sectional rank'),
       ('   B1 — B0 + static graph: degree centrality, Hype Hub Score, '
        'Louvain community one-hot'),
       ('   B2 — B1 + rolling attention (10-day mention count) + '
        'rolling sentiment (5-day mean)'),
       ('   B3 — XGBoost on B2 features (non-linear ceiling test)'),
   ],
   size=13, bullet=True)
notes(s, (
    'Path B setup — เน้น controls เพราะ finance data noisy\n\n'
    'Task: binary classification — up/down next day\n\n'
    'Expanding-window CV = train แต่ละ fold ยาวขึ้นเรื่อยๆ, test ถัดไป '
    '— realistic กว่า random split เพราะไม่มี lookahead\n\n'
    'Feature groups ทำเป็น ablation:\n'
    '- B0 = lag returns เท่านั้น (baseline)\n'
    '- B1 เพิ่ม static graph features\n'
    '- B2 เพิ่ม rolling attention + sentiment\n'
    '- B3 = XGBoost บน B2 (test ว่า non-linear help ไหม)\n\n'
    'Portfolio backtest: long-short 10% decile — practical relevance สำหรับ '
    'fund manager'
))


# ----- Slide 9: Path B Results ------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'สรุปผล · Path B Results', color='red')
title_of(s, 'Path B — Accuracy = flat, Sharpe = climbing 📈')
rows_txt = ['Model  Accuracy       95% CI            F1     Sharpe']
for k in ['B0', 'B1', 'B2', 'B3']:
    m = rp[k]
    rows_txt.append(
        f'{k:6s} {m["mean_accuracy"]:.3f}  '
        f'[{m["acc_ci_low"]:.3f}, {m["acc_ci_high"]:.3f}]   '
        f'{m["mean_f1"]:.3f}  {m["sharpe_proxy"]:>5.2f}'
    )
tb(s, Inches(0.7), Inches(1.4), Inches(7.5), Inches(2.5), rows_txt, size=13)
img(s, 'rp_sharpe_curve.png', Inches(8.4), Inches(1.3),
    width=Inches(4.7))
tb(s, Inches(0.7), Inches(4.4), Inches(12), Inches(2.8),
   [
       (f'Accuracy stays flat ~0.53–0.55 — SET50 direction ยากทำนายจริงๆ'),
       (f'Sharpe proxy climbs: B0 {b0["sharpe_proxy"]:.2f} → '
        f'B1 {b1["sharpe_proxy"]:.2f} → B2 {b2["sharpe_proxy"]:.2f}'),
       (f'B3 (XGBoost) Sharpe {b3["sharpe_proxy"]:.2f} — over-fit '
        '(non-linear ไม่ช่วยที่ N เล็ก)'),
       ('💡 Key insight: graph features moves the money, not the average — '
        'shift decision boundary ให้ portfolio ตอนเลือก "confident" pick '
        'แต่ไม่เพิ่ม accuracy overall'),
   ],
   size=13, bullet=True)
notes(s, (
    'Path B — ผลที่น่าสนใจที่สุด\n\n'
    'Accuracy ~0.55 flat ทั้ง 4 models — เกือบ chance level เพราะ SET50 '
    'direction ทำนายยาก (efficient market)\n\n'
    'แต่ Sharpe proxy โต monotonically: 3.07 → 3.79 → 4.30 ตอนใส่ graph + '
    'attention features\n\n'
    'ตีความ: features ไม่ได้ทำให้ทำนายถูก "โดยรวม" แต่ทำให้ตอนที่ model '
    'confident (top/bottom decile) ทำนายดีขึ้น → portfolio return ดีกว่า\n\n'
    'ข้อระวัง: Sharpe 4.30 สูงเกินจริงสำหรับ real trading (naive backtest, '
    'ไม่มี transaction cost, N=120 days) — ต้องระบุใน limitations\n\n'
    'XGBoost B3 แย่กว่า B2 — เพราะ overfits ใน small time-series sample'
))


# ----- Slide 10: Path C methodology + results ---------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'สรุปผล · Path C', color='red')
title_of(s, 'Path C — Fine-tune WangchanBERTa (Midterm limitation)')
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.0),
   [
       (f'Design: stratified 5-fold CV on 50 gold posts, 3 epochs/fold, '
        f'AdamW lr 2e-5, batch 4'),
   ],
   size=14, bullet=True)
rows_txt = ['Metric                    Baseline    Fine-tuned (5-fold)']
rows_txt.append(f'Accuracy                    {val["baseline"]["accuracy"]:.2f}          '
                f'{val["finetuned_pooled"]["accuracy"]:.2f} '
                f'(±{val["finetuned_cv_mean"]["accuracy_std"]:.2f})')
rows_txt.append(f'Macro F1                    {val["baseline"]["macro_f1"]:.2f}          '
                f'{val["finetuned_pooled"]["macro_f1"]:.2f}')
rows_txt.append(f'Cohen κ                     {val["baseline"]["kappa"]:.2f}          '
                f'{val["finetuned_pooled"]["kappa"]:.2f}')
rows_txt.append(f'Positive recall             {val["per_class_metrics"]["Positive"]["baseline"]["recall"]:.2f}          '
                f'{val["per_class_metrics"]["Positive"]["finetuned"]["recall"]:.2f}')
tb(s, Inches(0.7), Inches(2.6), Inches(7.5), Inches(2), rows_txt, size=14)
img(s, 'sentiment_finetuned_confusion.png', Inches(8.4), Inches(1.4),
    width=Inches(4.7))
tb(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2.3),
   [
       ('Accuracy 0.34 → 0.50 (+16pp, +47% relative) '
        '— above chance ระดับ 3-class (0.33)'),
       ('Positive recall 0.00 → 0.30 — model เริ่มทำนาย Positive class ได้'),
       ('Cohen κ 0.14 → 0.17 — improvement เล็ก (N=50 posts เล็ก)'),
       ('💡 H3 weakness = partly classifier artefact, ไม่ใช่ pure corpus '
        'property — เพิ่ม labels 200-500 → κ >0.5 ได้'),
   ],
   size=13, bullet=True)
notes(s, (
    'Path C — แก้จุดอ่อนใหญ่ของ Midterm\n\n'
    'ตอน Midterm อาจารย์ review บอกว่า classifier chance-level = จุดที่ต้องแก้\n\n'
    'ผลจริง: accuracy 0.34 → 0.50 — improvement ชัดเจน แต่ κ ยังต่ำ\n\n'
    'Positive recall เป็นตัวเด่น: baseline 0.00 (ไม่เคย predict Positive!) → '
    '0.30 หลัง fine-tune → เริ่มเห็น Positive class จริงๆ\n\n'
    'ข้อจำกัด: 50 gold posts เล็กเกินไปสำหรับ fine-tune BERT (105M params) — '
    'ถ้าเพิ่ม labels 200-500 posts κ น่าจะขึ้นถึง 0.5+\n\n'
    'Findings: H3 ที่อ่อนใน Midterm ส่วนหนึ่งเป็น classifier ceiling ไม่ใช่ '
    'corpus property — เป็น constructive finding ให้ future work'
))


# ----- Slide 11: Summary (3 findings) -----------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'สรุปผลรวม · Summary', color='red')
title_of(s, 'สรุป — 3 predictive stress tests, 3 findings')
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5),
   [
       ('🥇 Path A — Preferential Attachment ชนะ (AUC 0.73, P@10 = 1.00). '
        'Retail attention เป็น self-reinforcing → predict pair ที่ hub '
        'จะดัดกันเอง; overlooked co-movers ยังโดนละเลย'),
       (''),
       ('🥈 Path B — Graph features ไม่เพิ่ม accuracy (~0.55) แต่ **เพิ่ม '
        'Sharpe 3.07→4.30**. Attention moves the money, not the average. '
        'XGBoost overfits ที่ N=120 days'),
       (''),
       ('🥉 Path C — WangchanBERTa fine-tune บน 50 posts: accuracy 0.34→0.50, '
        'Positive recall 0.00→0.30. H3 weakness เป็น classifier artefact '
        'มากกว่า corpus property'),
       (''),
       ('🎯 Common thread: predictive extension ยืนยัน Midterm finding '
        'ที่ว่า "attention ≠ fundamentals" — retail hype self-reinforces, '
        'overlooked co-movers ยัง overlooked, แต่ hype signal ยัง usable '
        'ใน portfolio construction'),
   ],
   size=14)
notes(s, (
    'Slide สรุปหลัก\n\n'
    '3 findings ที่ survive predictive stress test:\n\n'
    '1. Path A: heuristics ชนะ + retail self-reinforce\n'
    '2. Path B: graph features usable for portfolio ไม่ใช่ accuracy\n'
    '3. Path C: H3 weakness partly classifier artefact\n\n'
    'Common thread ที่เชื่อม 3 paths กลับไปหา Midterm: attention ≠ '
    'fundamentals ยังเป็นจริง แต่ hype signal ยัง extractable ได้'
))


# ----- Slide 12: Course coverage ---------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'สรุปผลรวม · Course Coverage', color='red')
title_of(s, 'ครอบคลุมเทคนิคจาก Weeks 3–7')
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5),
   [
       ('Week 3 (Neo4j GDS + heuristics) — Jaccard, Adamic-Adar, '
        'Preferential Attachment ใน Path A'),
       ('Week 4 (Similarity + KNN) — Hadamard product z_u ⊙ z_v ใน '
        'Node2Vec classifier'),
       ('Week 5 (DeepWalk / Node2Vec / GCN) — Node2Vec 64-dim '
        '+ LogReg ใน Path A'),
       ('Week 6 (GNN concepts) — inductive bias discussion + GraphSAGE '
        'motivation'),
       ('Week 7 (Link Prediction on MovieLens) — Path A directly mirrors '
        'this workshop (HeteroData → SAGEConv → BCE + Adam)'),
       (''),
       ('Extra beyond weeks: Statistical inference (bootstrap CI, '
        'permutation), XGBoost, WangchanBERTa fine-tune, portfolio backtest'),
   ],
   size=14, bullet=True)
notes(s, (
    'Slide นี้ตอบโจทย์ "coverage" ของ course rubric\n\n'
    'Weeks 3-7 ทุก week มีการใช้จริงใน Final:\n'
    '- Week 3 heuristics = 3 methods\n'
    '- Week 4 similarity concept = Hadamard product\n'
    '- Week 5 Node2Vec = 1 method\n'
    '- Week 6 GNN concepts = discussed\n'
    '- Week 7 GraphSAGE = 1 method + directly mirrors MovieLens workshop\n\n'
    'Bonus: statistical inference (bootstrap CI), XGBoost, WangchanBERTa '
    'fine-tune, portfolio backtest — ไปเกินกว่าที่เรียน'
))


# ----- Slide 13: Limitations + Future -----------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
chip(s, 'แนวทางการพัฒนาต่อ · Future Work', color='green')
title_of(s, 'ข้อจำกัด + Future Work')
tb(s, Inches(0.7), Inches(1.4), Inches(12), Inches(2.5),
   [
       'ข้อจำกัดปัจจุบัน:',
       ('   Corpus 96 posts (event-anchored) — statistical power จำกัด'),
       ('   Trading window 120 days — เล็กสำหรับ Sharpe stability'),
       ('   Sentiment gold set 50 posts — เล็กสำหรับ fine-tune 105M-param model'),
       ('   Portfolio backtest naive (no transaction cost, no slippage)'),
   ],
   size=14, bullet=True, color='grey')
tb(s, Inches(0.7), Inches(4.3), Inches(12), Inches(3),
   [
       'Future work (เรียงตาม effort × impact):',
       ('   1. เพิ่ม gold labels 200-500 posts + fine-tune WangchanBERTa → '
        'H3 น่าจะ flip'),
       ('   2. Multi-event study — เลือก 3-5 events (BoT, earnings, index '
        'rebalance) test post-event density diffusion generalize ไหม'),
       ('   3. Aspect-based sentiment — สำหรับ post หลายหุ้น (DELTA บวก, '
        'OR ลบ ใน post เดียว)'),
       ('   4. Multi-platform — X + Facebook groups + Pantip triangulate'),
       ('   5. Dynamic Louvain — daily rolling communities แทน static'),
       ('   6. Realistic backtest — transaction cost, market impact, '
        'longer horizon'),
   ],
   size=13, bullet=True)
notes(s, (
    'Limitations honesty — สำคัญมากสำหรับ grade\n\n'
    'ข้อจำกัด: N เล็กในทุก dimension — corpus 96, trading 120 days, gold 50\n\n'
    'Future work เรียงตาม effort × impact:\n'
    '1. เพิ่ม labels ก่อน — quick win (1 คนใช้เวลา 2-3 ชม.)\n'
    '2. Multi-event study — ต้องขยาย window ยาวขึ้น\n'
    '3-6. เป็น long-term extensions\n\n'
    'ถ้าอาจารย์ถาม "what next" ตอบ direction 1 ก่อน — quick + impactful'
))


# ----- Slide 14: Thank you ----------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[5])
title_of(s, 'Thank you · Q&A', size=44)
tb(s, Inches(0.7), Inches(3.2), Inches(12), Inches(1.5),
   ['ทุก code + data + report reproducible ที่ Final/ (mirror Midterm/)'],
   size=18, color='grey', italic=True)
tb(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5),
   [
       'Repo: Final/',
       ('scripts/ (09_link_prediction, 10_return_prediction, '
        '11_sentiment_finetune, 12_report_final_latex, 13_make_slides)'),
       'report/ (DADS7201_Final_Report.tex + .pdf, DADS7201_Final_Slides.pptx)',
       'output/ (19 artefacts — CSVs, JSONs, PNG+SVG figures)',
   ],
   size=14, color='grey')
notes(s, (
    'ปิดการนำเสนอ\n\n'
    'ทุกอย่าง reproducible: rerun pipeline = ได้ตัวเลขเดิม\n\n'
    'Deliverables 3 อัน: 4-page LNCS report (.pdf + .tex), '
    'slide deck (.pptx), + 19 artefacts ใน output/\n\n'
    'พร้อมตอบคำถามครับ'
))


# ----- save --------------------------------------------------------------
out = REPORT_DIR / 'DADS7201_Final_Slides.pptx'
prs.save(str(out))
print(f'Saved {len(prs.slides)} slides to {out}')
print(f'  Speaker notes attached: {len(prs.slides)} slides')
